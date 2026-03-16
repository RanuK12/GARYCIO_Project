"""
Actualiza la presentación GARYCIO con los últimos cambios del proyecto.
Usa python-pptx para editar las diapositivas existentes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import os

INPUT = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "GARYCIO_Presentacion.pptx")
OUTPUT = INPUT  # overwrite

def find_and_replace_text(slide, old_text, new_text):
    """Replace text in all shapes of a slide, preserving formatting."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs)
            if old_text in full_text:
                # Replace in the first run that contains part of the text
                for run in paragraph.runs:
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                        return True
    return False


def update_slide_text_by_shape(slide, shape_idx, new_paragraphs):
    """
    Replace all text in a specific shape with new paragraphs.
    new_paragraphs is a list of (text, bold, size) tuples.
    """
    shapes = list(slide.shapes)
    if shape_idx >= len(shapes):
        return
    shape = shapes[shape_idx]
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear existing paragraphs (keep first, remove rest)
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    for i, (text, bold, size) in enumerate(new_paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
            p.clear()
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        if bold is not None:
            run.font.bold = bold
        if size is not None:
            run.font.size = Pt(size)


def list_slide_shapes(slide, slide_num):
    """Debug: list all shapes and their text in a slide."""
    print(f"\n=== Slide {slide_num} ===")
    for i, shape in enumerate(slide.shapes):
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text[:80].replace("\n", " | ")
        print(f"  [{i}] name='{shape.name}' text='{text}'.encode('ascii', 'replace').decode()")


def main():
    prs = Presentation(INPUT)

    # ── Slide 4 (Fase 1 - index 3): Update to include new features ──
    slide4 = prs.slides[3]

    # Update bullet items for Bot de WhatsApp section
    replacements_slide4 = [
        ("Contacto automatico con donantes de zonas nuevas",
         "Contacto automatico con donantes: confirma si donan, dias y direccion exacta"),
        ("Confirmacion de dias de recoleccion disponibles",
         "Recoleccion de direccion exacta para reorganizar recorridos"),
        ("Presentacion del servicio y recoleccion de datos",
         "Registro de datos y presentacion del servicio"),
    ]

    for old, new in replacements_slide4:
        find_and_replace_text(slide4, old, new)

    # Add incident and report info to Fase 1
    replacements_reclamos = [
        ("Canal para avisos de donantes (no recoleccion, etc.)",
         "Reporte de incidentes con alerta inmediata al CEO"),
        ("Notificaciones automaticas al equipo",
         "Reporte PDF diario a las 19:00 hs con estadisticas"),
        ("Registro y seguimiento de cada caso",
         "Registro de litros, bidones y combustible por chofer"),
    ]
    for old, new in replacements_reclamos:
        find_and_replace_text(slide4, old, new)

    # ── Slide 8 (Que necesitamos - index 7): Update requirements ──
    slide8 = prs.slides[7]

    replacements_slide8 = [
        ("Numeros de telefono, direcciones, zonas asignadas y estado actual de cada donante.",
         "Numeros de telefono, direcciones EXACTAS (calle, numero, entre calles, piso), zonas y estado de donacion actual."),
        ("Detalle de las zonas recientemente asignadas: limites, donantes potenciales, contactos.",
         "Detalle de zonas nuevas: limites, donantes potenciales. Confirmar si estan donando y que dia."),
    ]
    for old, new in replacements_slide8:
        find_and_replace_text(slide8, old, new)

    # ── Slide 9 (Cronograma - index 8): Update timeline ──
    slide9 = prs.slides[8]

    timeline_updates = [
        ("Setup + Bot WhatsApp Basico", "Bot WhatsApp + Choferes"),
        ("Configuracion de infraestructura", "Bot con 7 flujos operativos"),
        ("Bot WhatsApp: flujo de presentacion", "Sistema de choferes e incidentes"),
        ("Integracion con WhatsApp Business API", "Reportes PDF diarios automaticos"),
        ("Envio masivo personalizado a donantes", "Contacto masivo a nuevas zonas"),
        ("Recoleccion de datos (dias, estado)", "Recoleccion de direccion exacta"),
        ("Base de datos de donantes activos", "Base de datos + migracion datos"),
    ]
    for old, new in timeline_updates:
        find_and_replace_text(slide9, old, new)

    # ── Save ──
    prs.save(OUTPUT)
    print(f"\nPresentacion actualizada: {OUTPUT}")
    size_kb = os.path.getsize(OUTPUT) / 1024
    print(f"Tamano: {size_kb:.1f} KB")


if __name__ == "__main__":
    main()
